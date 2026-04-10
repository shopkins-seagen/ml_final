from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG   = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT    = RGBColor(0x58, 0x9C, 0xF4)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xFF)
SUBTEXT   = RGBColor(0xAA, 0xBB, 0xDD)
GREEN     = RGBColor(0x4C, 0xAF, 0x50)
ORANGE    = RGBColor(0xFF, 0x99, 0x22)

BLANK = prs.slide_layouts[6]   # completely blank


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, text, font_size=18, bold=False,
        color=WHITE, bg_color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox


def rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def accent_bar(slide, w=0.07):
    rect(slide, 0, 0, w, 7.5, ACCENT)


# ── Slide 1 : Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
rect(s, 0.07, 3.1, 13.26, 0.06, ACCENT)
box(s, 0.5, 0.8, 12, 1.2,
    "SDTM Domain Classifier",
    font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 0.5, 2.1, 12, 0.7,
    "Classifying Clinical Trial Variables using KNN & CART",
    font_size=22, color=LIGHT, align=PP_ALIGN.CENTER)
box(s, 0.5, 3.4, 12, 0.6,
    "Domains Covered:  AE · CM · DM · EG · EX · LB · MH · PE · SC · VS",
    font_size=17, color=SUBTEXT, align=PP_ALIGN.CENTER)
box(s, 0.5, 4.2, 12, 0.5,
    "Models:  K-Nearest Neighbors (KNN)   vs   Classification & Regression Tree (CART)",
    font_size=16, color=SUBTEXT, align=PP_ALIGN.CENTER)
box(s, 0.5, 5.2, 12, 0.5,
    "Dataset: 266 SDTM variables  ·  200 TF-IDF features  ·  10 target classes",
    font_size=14, italic=True, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ── Slide 2 : Project Overview ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.3, 12, 0.7, "Project Overview", font_size=30, bold=True, color=ACCENT)
rect(s, 0.5, 1.1, 12, 0.05, ACCENT)

items = [
    ("Objective",
     "Given a variable label from a clinical trial dataset, automatically predict\nwhich SDTM domain it belongs to (AE, LB, VS, etc.)."),
    ("Why it matters",
     "Manual SDTM mapping is time-consuming and error-prone.\nAn ML classifier speeds up dataset curation for clinical submissions."),
    ("Approach",
     "Text features (TF-IDF) from variable names & labels → train KNN and CART → compare performance."),
    ("Dataset",
     "266 SDTM variables across 10 domains, sourced from CDISC SDTM Implementation Guide."),
]
y = 1.3
for title, desc in items:
    rect(s, 0.5, y, 3.2, 0.9, RGBColor(0x25, 0x25, 0x45))
    box(s, 0.6, y + 0.05, 3.0, 0.35, title, font_size=13, bold=True, color=ACCENT)
    box(s, 0.6, y + 0.38, 3.0, 0.55, desc, font_size=10, color=WHITE)
    y += 1.1

# ── Slide 3 : Pipeline Overview ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.3, 12, 0.7, "End-to-End Pipeline", font_size=30, bold=True, color=ACCENT)
rect(s, 0.5, 1.1, 12, 0.05, ACCENT)

steps = [
    ("1", "Load Data",        "Read sdtm_variables.csv\n266 rows × 3 cols"),
    ("2", "Preprocess",       "Drop nulls\nLabel-encode domain"),
    ("3", "Feature Eng.",     "Concat name+label\nTF-IDF (200 features)"),
    ("4", "Split",            "75% train / 25% test\nStratified by domain"),
    ("5", "Train KNN",        "CV over k=1..15\nCosine distance"),
    ("6", "Train CART",       "GridSearch depth\n& min_samples_split"),
    ("7", "Evaluate",         "Accuracy, F1\nConfusion matrix"),
    ("8", "Compare",          "KNN vs CART\nBar chart summary"),
]
x = 0.5
for num, title, desc in steps:
    rect(s, x, 1.5, 1.45, 1.5, RGBColor(0x25, 0x25, 0x45))
    box(s, x, 1.52, 1.45, 0.45, num, font_size=22, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)
    box(s, x, 1.95, 1.45, 0.4, title, font_size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    box(s, x, 2.35, 1.45, 0.65, desc, font_size=9, color=SUBTEXT,
        align=PP_ALIGN.CENTER)
    if x < 11.5:
        box(s, x + 1.45, 2.1, 0.25, 0.35, "→", font_size=16,
            color=ACCENT, align=PP_ALIGN.CENTER)
    x += 1.6

# ── Slide 4 : Step 1 — Load Data ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 1 — Load & Explore the Dataset",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.15, 5.5, 0.4, "Dataset: sdtm_variables.csv",
    font_size=15, bold=True, color=WHITE)
cols = [("variable_name", "AETERM"), ("variable_label", "Reported Term for the Adverse Event"),
        ("domain", "AE")]
y = 1.65
for col, ex in cols:
    box(s, 0.5, y, 2.5, 0.38, col, font_size=12, bold=True, color=ACCENT)
    box(s, 3.1, y, 3.5, 0.38, ex, font_size=12, color=LIGHT)
    y += 0.42

box(s, 0.5, 3.1, 5.5, 0.4, "Domain Distribution (266 rows):",
    font_size=14, bold=True, color=WHITE)
domains = [("AE", 28), ("CM", 28), ("DM", 32), ("EG", 28), ("EX", 23),
           ("LB", 31), ("MH", 24), ("PE", 22), ("SC", 16), ("VS", 34)]
x = 0.5
for dom, cnt in domains:
    rect(s, x, 3.6, 1.15, 1.2 * cnt / 34, RGBColor(0x58, 0x9C, 0xF4))
    box(s, x, 3.6 + 1.2 * cnt / 34, 1.15, 0.35,
        f"{dom}\n{cnt}", font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
    x += 1.2

box(s, 7.0, 1.15, 6.0, 5.5,
    "Key observations:\n\n"
    "• 3 columns: variable_name, variable_label, domain\n\n"
    "• 10 SDTM domains are the target classes\n\n"
    "• No missing values in the dataset\n\n"
    "• Domain counts range from 16 (SC) to 34 (VS)\n\n"
    "• Original: 198 rows\n  After augmentation: 266 rows (+68 synthetic)",
    font_size=13, color=WHITE)

# ── Slide 5 : Step 2 — Preprocessing ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 2 — Data Preprocessing & Label Encoding",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.2, 5.8, 0.4, "2a — Drop missing values (defensive)",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 1.65, 5.8, 0.5,
    "df_clean = df.dropna().copy()\n→ 266 clean records", font_size=12, color=LIGHT,
    bg_color=RGBColor(0x20, 0x20, 0x38))

box(s, 0.5, 2.4, 5.8, 0.4, "2b — Label Encode the domain column",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 2.85, 5.8, 0.5,
    "le = LabelEncoder()\ndf_clean['domain_encoded'] = le.fit_transform(df_clean['domain'])",
    font_size=12, color=LIGHT, bg_color=RGBColor(0x20, 0x20, 0x38))

box(s, 0.5, 3.55, 5.8, 0.4, "Why label encoding?",
    font_size=13, bold=True, color=ACCENT)
box(s, 0.5, 3.95, 5.8, 0.8,
    "ML models need numeric targets, not strings.\n"
    "LabelEncoder maps each domain to an integer (0–9).",
    font_size=12, color=WHITE)

# mapping table
box(s, 7.0, 1.2, 5.8, 0.4, "Domain → Encoded Label Mapping",
    font_size=14, bold=True, color=WHITE)
mapping = [("AE",0),("CM",1),("DM",2),("EG",3),("EX",4),
           ("LB",5),("MH",6),("PE",7),("SC",8),("VS",9)]
y = 1.7
for dom, code in mapping:
    rect(s, 7.0, y, 2.0, 0.38, RGBColor(0x25, 0x25, 0x45))
    box(s, 7.0, y, 1.0, 0.38, dom, font_size=13, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)
    box(s, 8.1, y, 0.9, 0.38, f"→  {code}", font_size=13, color=WHITE)
    y += 0.42

# ── Slide 6 : Step 3 — Feature Engineering ───────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 3 — Feature Engineering (TF-IDF)",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.15, 12, 0.4, "Combine variable_name + variable_label into one text string per row:",
    font_size=13, color=WHITE)
box(s, 0.5, 1.6, 12, 0.55,
    "df_clean['text'] = df_clean['variable_name'].str.lower() + ' ' + df_clean['variable_label'].str.lower()",
    font_size=11, color=LIGHT, bg_color=RGBColor(0x20, 0x20, 0x38))

box(s, 0.5, 2.35, 5.8, 0.4, "TF-IDF Parameters", font_size=14, bold=True, color=WHITE)
params = [
    ("ngram_range=(1,2)", "Captures unigrams AND bigrams\ne.g. 'adverse' + 'adverse event'"),
    ("max_features=200",  "Keeps top 200 most informative tokens\nout of 1000+ possible"),
    ("sublinear_tf=True", "Uses log(TF)+1 instead of raw TF\nReduces dominance of frequent words"),
    ("stop_words='english'","Removes 'the','for','of','in'\nReduces noise"),
]
y = 2.85
for param, desc in params:
    rect(s, 0.5, y, 5.8, 0.68, RGBColor(0x25, 0x25, 0x45))
    box(s, 0.6, y + 0.02, 2.6, 0.3, param, font_size=11, bold=True, color=ORANGE)
    box(s, 0.6, y + 0.32, 5.5, 0.35, desc, font_size=10, color=WHITE)
    y += 0.75

box(s, 7.0, 2.35, 6.0, 0.4, "Output: Feature Matrix X", font_size=14, bold=True, color=WHITE)
box(s, 7.0, 2.85, 6.0, 1.8,
    "Shape:  (266, 200)\n\n"
    "• 266 rows = one per SDTM variable\n"
    "• 200 cols = one per TF-IDF token\n"
    "• Values = TF-IDF score (0.0 if absent)\n\n"
    "Example row for 'aeterm adverse event':\n"
    "  adverse       →  0.82\n"
    "  adverse event →  0.71\n"
    "  laboratory    →  0.00\n"
    "  vital signs   →  0.00",
    font_size=12, color=WHITE, bg_color=RGBColor(0x20, 0x20, 0x38))

# ── Slide 7 : Step 4 — Train/Test Split ──────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 4 — Train / Test Split",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.2, 12, 0.45,
    "train_test_split(X, y, test_size=0.25, random_state=42, stratify=y)",
    font_size=13, color=LIGHT, bg_color=RGBColor(0x20, 0x20, 0x38))

# visual split bar
rect(s, 0.5, 1.9, 9.0, 0.9, GREEN)
rect(s, 9.5, 1.9, 3.0, 0.9, ORANGE)
box(s, 0.5, 1.9, 9.0, 0.9, "TRAINING SET\n148 samples (75%)",
    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 9.5, 1.9, 3.0, 0.9, "TEST SET\n50 samples (25%)",
    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.5, 3.05, 5.8, 0.4, "Key parameters:", font_size=13, bold=True, color=WHITE)
notes = [
    ("test_size=0.25",   "25% held out for final evaluation"),
    ("random_state=42",  "Reproducible split every run"),
    ("stratify=y",       "Each domain keeps same proportion\nin both train and test sets"),
]
y = 3.5
for p, d in notes:
    box(s, 0.5, y, 2.5, 0.38, p, font_size=12, bold=True, color=ACCENT)
    box(s, 3.1, y, 3.2, 0.38, d, font_size=12, color=WHITE)
    y += 0.48

box(s, 7.0, 3.05, 5.8, 0.4, "Why stratify?", font_size=13, bold=True, color=WHITE)
box(s, 7.0, 3.5, 5.8, 1.5,
    "Without stratify, a domain with only 16 rows (SC)\n"
    "could end up with 0 samples in the test set.\n\n"
    "stratify=y guarantees every domain appears\n"
    "in both splits in its correct proportion.",
    font_size=12, color=WHITE)

# ── Slide 8 : Step 5 — KNN ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 5 — Train KNN Model",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.15, 5.8, 0.4, "How KNN works", font_size=14, bold=True, color=WHITE)
box(s, 0.5, 1.6, 5.8, 1.9,
    "1. Memorizes all 148 training samples\n"
    "2. Given a test sample, computes COSINE\n"
    "   DISTANCE to every training point\n"
    "3. Finds k nearest neighbors\n"
    "4. Takes majority vote → predicted domain\n\n"
    "Why cosine? TF-IDF vectors are sparse &\n"
    "high-dimensional. Cosine measures direction\n"
    "(word usage pattern), not magnitude.",
    font_size=12, color=WHITE)

box(s, 0.5, 3.7, 5.8, 0.4, "Hyperparameter tuning: find best k",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 4.15, 5.8, 1.0,
    "for k in range(1, 16):\n"
    "    5-fold CV on X_train\n"
    "    record mean accuracy\n"
    "best_k = k with highest CV accuracy",
    font_size=12, color=LIGHT, bg_color=RGBColor(0x20, 0x20, 0x38))

box(s, 7.0, 1.15, 5.8, 0.4, "Cross-Validation explained",
    font_size=14, bold=True, color=WHITE)
box(s, 7.0, 1.6, 5.8, 2.5,
    "5-fold CV splits training data into 5 parts:\n\n"
    "  Fold 1: [Val] [Tr] [Tr] [Tr] [Tr]\n"
    "  Fold 2: [Tr] [Val] [Tr] [Tr] [Tr]\n"
    "  Fold 3: [Tr] [Tr] [Val] [Tr] [Tr]\n"
    "  Fold 4: [Tr] [Tr] [Tr] [Val] [Tr]\n"
    "  Fold 5: [Tr] [Tr] [Tr] [Tr] [Val]\n\n"
    "Average accuracy across 5 folds = CV score\n"
    "Test data is NEVER seen during this step.",
    font_size=12, color=WHITE)

# ── Slide 9 : Step 6 — CART ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 6 — Train CART Model",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.15, 5.8, 0.4, "How CART works", font_size=14, bold=True, color=WHITE)
box(s, 0.5, 1.6, 5.8, 2.0,
    "Builds a binary decision tree by repeatedly\n"
    "splitting the data on the feature that most\n"
    "reduces GINI IMPURITY:\n\n"
    "  Gini = 1 - Σ(pᵢ)²\n\n"
    "  • Gini = 0  → perfectly pure node\n"
    "  • Gini = 0.9 → very mixed node\n\n"
    "At each node: 'Is token X score > 0.3?'\n"
    "→ YES branch / NO branch → leaf = domain",
    font_size=12, color=WHITE)

box(s, 0.5, 3.8, 5.8, 0.4, "GridSearchCV parameter grid",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 4.25, 5.8, 0.9,
    "max_depth: [3, 5, 7, 10, None]\n"
    "min_samples_split: [2, 4, 6]\n"
    "→ 5 × 3 = 15 combinations × 5-fold CV = 75 fits",
    font_size=12, color=LIGHT, bg_color=RGBColor(0x20, 0x20, 0x38))

box(s, 7.0, 1.15, 5.8, 0.4, "Key hyperparameters",
    font_size=14, bold=True, color=WHITE)
cart_params = [
    ("max_depth",         "Limits tree depth → prevents overfitting\nNone = grow until pure leaves"),
    ("min_samples_split", "Min samples needed to split a node\nHigher = simpler tree"),
    ("criterion='gini'",  "Splitting criterion\nGini impurity measures node purity"),
]
y = 1.6
for p, d in cart_params:
    rect(s, 7.0, y, 5.8, 0.78, RGBColor(0x25, 0x25, 0x45))
    box(s, 7.1, y + 0.02, 5.6, 0.3, p, font_size=12, bold=True, color=ORANGE)
    box(s, 7.1, y + 0.34, 5.6, 0.42, d, font_size=11, color=WHITE)
    y += 0.85

# ── Slide 10 : Evaluation Metrics ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Steps 7 & 8 — Evaluation & Comparison",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

metrics = [
    ("Accuracy",  "% of all predictions correct\n(TP+TN) / Total"),
    ("Precision", "Of predicted domain X,\nhow many were actually X?"),
    ("Recall",    "Of actual domain X,\nhow many were found?"),
    ("F1-Score",  "Harmonic mean of Precision & Recall\nBalances both metrics"),
]
x = 0.5
for m, d in metrics:
    rect(s, x, 1.2, 2.9, 1.3, RGBColor(0x25, 0x25, 0x45))
    box(s, x, 1.22, 2.9, 0.45, m, font_size=15, bold=True,
        color=ACCENT, align=PP_ALIGN.CENTER)
    box(s, x, 1.7, 2.9, 0.75, d, font_size=11, color=WHITE,
        align=PP_ALIGN.CENTER)
    x += 3.1

box(s, 0.5, 2.75, 5.8, 0.4, "Confusion Matrix", font_size=14, bold=True, color=WHITE)
box(s, 0.5, 3.2, 5.8, 1.5,
    "A 10×10 matrix showing:\n"
    "  • Rows = actual domain\n"
    "  • Cols = predicted domain\n"
    "  • Diagonal = correct predictions\n"
    "  • Off-diagonal = misclassifications\n"
    "Helps identify which domains get confused.",
    font_size=12, color=WHITE)

box(s, 7.0, 2.75, 5.8, 0.4, "Weighted averaging",
    font_size=14, bold=True, color=WHITE)
box(s, 7.0, 3.2, 5.8, 1.5,
    "Metrics are computed per-class then\n"
    "averaged weighted by class size.\n\n"
    "Important because domains are slightly\n"
    "imbalanced (SC=16 vs VS=34).\n"
    "Weighted avg prevents larger domains\n"
    "from masking poor performance on smaller ones.",
    font_size=12, color=WHITE)

# ── Slide 11 : CART Interpretability ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Step 9 — Visualize CART Tree & Feature Importance",
    font_size=28, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

box(s, 0.5, 1.15, 5.8, 0.4, "Decision Tree Visualization",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 1.6, 5.8, 1.5,
    "Top 3 levels of the tree are plotted.\n"
    "Each node shows:\n"
    "  • Splitting feature (TF-IDF token)\n"
    "  • Threshold value\n"
    "  • Gini impurity\n"
    "  • Sample count\n"
    "  • Predicted class (colour-coded)",
    font_size=12, color=WHITE)

box(s, 0.5, 3.3, 5.8, 0.4, "Feature Importances",
    font_size=14, bold=True, color=WHITE)
box(s, 0.5, 3.75, 5.8, 1.6,
    "importance = total Gini reduction from splits\non that feature across the whole tree.\n\n"
    "Top features reveal which tokens are\nmost discriminative, e.g.:\n"
    "  'adverse'   → strongly predicts AE\n"
    "  'laboratory'→ strongly predicts LB\n"
    "  'vital'     → strongly predicts VS",
    font_size=12, color=WHITE)

box(s, 7.0, 1.15, 5.8, 0.4, "Why interpretability matters in clinical context",
    font_size=14, bold=True, color=WHITE)
box(s, 7.0, 1.6, 5.8, 3.5,
    "SDTM mapping is a regulated process.\n"
    "Auditors and data managers need to understand\n"
    "WHY a variable was classified into a domain.\n\n"
    "CART provides a clear audit trail:\n"
    "  'If adverse event > 0.3 → predict AE'\n\n"
    "KNN gives no such explanation — it only\n"
    "says 'the 3 most similar variables in\n"
    "training were all AE'.\n\n"
    "→ CART is more trustworthy in clinical settings\n"
    "  even if KNN has slightly higher accuracy.",
    font_size=12, color=WHITE)

# ── Slide 12 : Summary ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s)
box(s, 0.5, 0.25, 12, 0.65, "Summary & Conclusions",
    font_size=30, bold=True, color=ACCENT)
rect(s, 0.5, 1.0, 12, 0.05, ACCENT)

cols_data = [
    ("KNN", ACCENT, [
        "Instance-based (lazy learner)",
        "Stores all training data",
        "Cosine distance on TF-IDF",
        "Best k chosen by CV",
        "No explicit decision rules",
        "Black-box predictions",
    ]),
    ("CART", GREEN, [
        "Rule-based tree (eager learner)",
        "Learns explicit split rules",
        "Gini impurity criterion",
        "max_depth tuned by GridSearch",
        "Fully interpretable tree",
        "Feature importance scores",
    ]),
]
x = 0.5
for title, col, points in cols_data:
    rect(s, x, 1.2, 5.8, 0.55, col)
    box(s, x, 1.2, 5.8, 0.55, title, font_size=20, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    y = 1.85
    for pt in points:
        box(s, x + 0.2, y, 5.5, 0.42, f"• {pt}", font_size=12, color=WHITE)
        y += 0.44
    x += 6.3

box(s, 0.5, 5.3, 12, 0.4, "Key Takeaway:", font_size=14, bold=True, color=ACCENT)
box(s, 0.5, 5.75, 12, 0.7,
    "Both models leverage TF-IDF features from SDTM variable names & labels. "
    "KNN benefits from cosine similarity on sparse text vectors, while CART "
    "provides interpretable rules — critical for clinical/regulatory use cases.",
    font_size=12, color=WHITE)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "SDTM_Domain_Classifier_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
