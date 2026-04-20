
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation(r"C:\dev\stevens\data_analytics\project\ML_domain_classifier-\template.pptx")
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1E,0x1E,0x2E); ACCENT = RGBColor(0x58,0x9C,0xF4)
WHITE   = RGBColor(0xFF,0xFF,0xFF);  LIGHT  = RGBColor(0xCC,0xDD,0xFF)
SUBTEXT = RGBColor(0xAA,0xBB,0xDD); GREEN  = RGBColor(0x4C,0xAF,0x50)
ORANGE  = RGBColor(0xFF,0x99,0x22);  RED    = RGBColor(0xEF,0x53,0x50)
PURPLE  = RGBColor(0xAB,0x47,0xBC)
BLANK   = prs.slide_layouts[0]

def bg(slide, color=DARK_BG):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, font_size=18, bold=False, color=WHITE,
        bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run()
    r.text = text; r.font.size = Pt(font_size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if bg_color:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg_color
    return tb

def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def accent_bar(slide): rect(slide, 0, 0, 0.07, 7.5, ACCENT)

def header(slide, title):
    accent_bar(slide)
    box(slide, 0.5, 0.25, 12, 0.65, title, 28, True, ACCENT)
    rect(slide, 0.5, 1.0, 12, 0.05, ACCENT)

def stat_tile(slide, l, t, w, h, value, label, vc=ACCENT):
    rect(slide, l, t, w, h, RGBColor(0x25,0x25,0x45))
    box(slide, l, t+0.08, w, h*0.52, value, 26, True, vc, align=PP_ALIGN.CENTER)
    box(slide, l, t+h*0.55, w, h*0.38, label, 10, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); accent_bar(s)
rect(s, 0.07, 3.1, 13.26, 0.06, ACCENT)
box(s,0.5,0.8,12,1.2,"SDTM Domain Classifier",44,True,WHITE,align=PP_ALIGN.CENTER)
box(s,0.5,2.1,12,0.7,"Classifying Clinical Trial Variables using KNN & CART",22,color=LIGHT,align=PP_ALIGN.CENTER)
box(s,0.5,3.4,12,0.6,"Domains: AE \u00b7 CM \u00b7 DM \u00b7 EG \u00b7 EX \u00b7 LB \u00b7 MH \u00b7 PE \u00b7 SC \u00b7 VS",17,color=SUBTEXT,align=PP_ALIGN.CENTER)
box(s,0.5,4.2,12,0.5,"Models:  K-Nearest Neighbors (KNN)   vs   Classification & Regression Tree (CART)",16,color=SUBTEXT,align=PP_ALIGN.CENTER)
box(s,0.5,4.9,12,0.5,"Dataset: 266 SDTM variables  \u00b7  200 TF-IDF features  \u00b7  10 target classes",14,italic=True,color=SUBTEXT,align=PP_ALIGN.CENTER)
box(s,0.5,6.3,12,0.45,"Data Analytics Final Project  \u00b7  Stevens Institute of Technology",13,italic=True,color=SUBTEXT,align=PP_ALIGN.CENTER)

# ── Slide 2: Project Overview ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Project Overview")
items = [
    ("Objective","Given a variable label from a clinical trial dataset, automatically predict which SDTM domain it belongs to (AE, LB, VS, etc.)."),
    ("Why it matters","Manual SDTM mapping is time-consuming and error-prone. An ML classifier speeds up dataset curation for clinical submissions."),
    ("Approach","Text features (TF-IDF) from variable names & labels \u2192 train KNN and CART \u2192 compare performance on held-out test set."),
    ("Dataset","266 SDTM variables across 10 domains, sourced from the CDISC SDTM Implementation Guide. Stratified 75/25 train/test split."),
]
y=1.3
for title,desc in items:
    rect(s,0.5,y,12.0,1.05,RGBColor(0x25,0x25,0x45))
    box(s,0.7,y+0.05,11.5,0.35,title,13,True,ACCENT)
    box(s,0.7,y+0.4,11.5,0.6,desc,12,color=WHITE)
    y+=1.15

# ── Slide 3: Pipeline ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"End-to-End Pipeline")
steps=[("1","Load Data","Read CSV\n266 rows"),("2","Preprocess","Drop nulls\nLabel encode"),
       ("3","TF-IDF","Vectorize\n200 features"),("4","Split","75% train\n25% test"),
       ("5","Train KNN","CV k=1..15\nCosine dist"),("6","Train CART","GridSearch\ndepth+split"),
       ("7","Evaluate","Accuracy\nF1 + CM"),("8","Compare","KNN vs CART\nBar chart")]
x=0.5
for num,title,desc in steps:
    rect(s,x,1.5,1.45,1.6,RGBColor(0x25,0x25,0x45))
    box(s,x,1.52,1.45,0.45,num,22,True,ACCENT,align=PP_ALIGN.CENTER)
    box(s,x,1.95,1.45,0.4,title,12,True,WHITE,align=PP_ALIGN.CENTER)
    box(s,x,2.35,1.45,0.7,desc,9,color=SUBTEXT,align=PP_ALIGN.CENTER)
    if x<11.5: box(s,x+1.45,2.1,0.25,0.35,"\u2192",16,color=ACCENT,align=PP_ALIGN.CENTER)
    x+=1.6

# ── Slide 4: Load Data ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 1 \u2014 Load & Explore the Dataset")
box(s,0.5,1.15,5.5,0.4,"Dataset: sdtm_variables.csv",15,True,WHITE)
for i,(col,ex) in enumerate([("variable_name","AETERM"),("variable_label","Reported Term for the Adverse Event"),("domain","AE")]):
    box(s,0.5,1.65+i*0.42,2.5,0.38,col,12,True,ACCENT)
    box(s,3.1,1.65+i*0.42,3.5,0.38,ex,12,color=LIGHT)
box(s,0.5,3.0,5.5,0.4,"Domain Distribution (266 rows):",14,True,WHITE)
domains=[("AE",28),("CM",28),("DM",32),("EG",28),("EX",23),("LB",31),("MH",24),("PE",22),("SC",16),("VS",34)]
x=0.5
for dom,cnt in domains:
    rect(s,x,3.5,1.1,1.2*cnt/34,ACCENT)
    box(s,x,3.5+1.2*cnt/34,1.1,0.38,f"{dom}\n{cnt}",9,color=WHITE,align=PP_ALIGN.CENTER)
    x+=1.15
box(s,7.0,1.15,6.0,5.0,
    "Key observations:\n\n\u2022 3 columns: variable_name, variable_label, domain\n\n"
    "\u2022 10 SDTM domains are the target classes\n\n\u2022 No missing values after cleaning\n\n"
    "\u2022 Domain counts range from 16 (SC) to 34 (VS)\n\n\u2022 Slight class imbalance \u2014 handled with stratified split",13,color=WHITE)

# ── Slide 5: Preprocessing ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 2 \u2014 Data Preprocessing & Label Encoding")
box(s,0.5,1.2,5.8,0.4,"2a \u2014 Drop missing values",14,True,WHITE)
box(s,0.5,1.65,5.8,0.5,"df_clean = df.dropna().copy()\n\u2192 266 clean records",12,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
box(s,0.5,2.4,5.8,0.4,"2b \u2014 Label Encode the domain column",14,True,WHITE)
box(s,0.5,2.85,5.8,0.5,"le = LabelEncoder()\ndf_clean['domain_encoded'] = le.fit_transform(df_clean['domain'])",12,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
box(s,0.5,3.55,5.8,0.4,"Why label encoding?",13,True,ACCENT)
box(s,0.5,3.95,5.8,0.8,"ML models need numeric targets, not strings.\nLabelEncoder maps each domain to an integer (0\u20139).",12,color=WHITE)
box(s,7.0,1.2,5.8,0.4,"Domain \u2192 Encoded Label Mapping",14,True,WHITE)
mapping=[("AE",0),("CM",1),("DM",2),("EG",3),("EX",4),("LB",5),("MH",6),("PE",7),("SC",8),("VS",9)]
y=1.7
for dom,code in mapping:
    rect(s,7.0,y,2.0,0.38,RGBColor(0x25,0x25,0x45))
    box(s,7.0,y,1.0,0.38,dom,13,True,ACCENT,align=PP_ALIGN.CENTER)
    box(s,8.1,y,0.9,0.38,f"\u2192  {code}",13,color=WHITE); y+=0.42

# ── Slide 6: Feature Engineering ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 3 \u2014 Feature Engineering (TF-IDF)")
box(s,0.5,1.15,12,0.4,"Combine variable_name + variable_label into one text string per row:",13,color=WHITE)
box(s,0.5,1.6,12,0.55,"df_clean['text'] = df_clean['variable_name'].str.lower() + ' ' + df_clean['variable_label'].str.lower()",11,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
box(s,0.5,2.35,5.8,0.4,"TF-IDF Parameters",14,True,WHITE)
params=[("ngram_range=(1,2)","Captures unigrams AND bigrams\ne.g. 'adverse' + 'adverse event'"),
        ("max_features=200","Keeps top 200 most informative tokens"),
        ("sublinear_tf=True","Uses log(TF)+1 instead of raw TF\nReduces dominance of frequent words"),
        ("stop_words='english'","Removes 'the','for','of','in' \u2014 reduces noise")]
y=2.85
for param,desc in params:
    rect(s,0.5,y,5.8,0.68,RGBColor(0x25,0x25,0x45))
    box(s,0.6,y+0.02,2.6,0.3,param,11,True,ORANGE); box(s,0.6,y+0.32,5.5,0.35,desc,10,color=WHITE); y+=0.75
box(s,7.0,2.35,6.0,0.4,"Output: Feature Matrix X",14,True,WHITE)
box(s,7.0,2.85,6.0,1.8,"Shape:  (266, 200)\n\n\u2022 266 rows = one per SDTM variable\n\u2022 200 cols = one per TF-IDF token\n\u2022 Values = TF-IDF score (0.0 if absent)\n\nExample row for 'aeterm adverse event':\n  adverse       \u2192  0.82\n  adverse event \u2192  0.71\n  vital signs   \u2192  0.00",12,color=WHITE,bg_color=RGBColor(0x20,0x20,0x38))

# ── Slide 7: Train/Test Split ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 4 \u2014 Train / Test Split")
box(s,0.5,1.2,12,0.45,"train_test_split(X, y, test_size=0.25, random_state=42, stratify=y)",13,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
rect(s,0.5,1.9,9.0,0.9,GREEN); rect(s,9.5,1.9,3.0,0.9,ORANGE)
box(s,0.5,1.9,9.0,0.9,"TRAINING SET\n198 samples (75%)",14,True,WHITE,align=PP_ALIGN.CENTER)
box(s,9.5,1.9,3.0,0.9,"TEST SET\n68 samples (25%)",14,True,WHITE,align=PP_ALIGN.CENTER)
box(s,0.5,3.05,5.8,0.4,"Key parameters:",13,True,WHITE)
for y,(p,d) in enumerate([("test_size=0.25","25% held out for final evaluation"),("random_state=42","Reproducible split every run"),("stratify=y","Each domain keeps same proportion in both splits")],0):
    box(s,0.5,3.5+y*0.52,2.5,0.38,p,12,True,ACCENT); box(s,3.1,3.5+y*0.52,3.2,0.38,d,12,color=WHITE)
box(s,7.0,3.05,5.8,0.4,"Why stratify?",13,True,WHITE)
box(s,7.0,3.5,5.8,1.5,"Without stratify, SC (only 16 rows) could end\nup with 0 samples in the test set.\n\nstratify=y guarantees every domain appears\nin both splits in its correct proportion.",12,color=WHITE)

# ── Slide 8: Train KNN ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 5 \u2014 Train KNN Model")
box(s,0.5,1.15,5.8,0.4,"How KNN works",14,True,WHITE)
box(s,0.5,1.6,5.8,2.0,"1. Memorizes all training samples\n2. Given a test point: computes COSINE\n   DISTANCE to every training point\n3. Finds k nearest neighbors\n4. Takes majority vote \u2192 predicted domain\n\nWhy cosine? TF-IDF vectors are sparse &\nhigh-dimensional. Cosine measures direction\n(word usage pattern), not magnitude.",12,color=WHITE)
box(s,0.5,3.7,5.8,0.4,"Hyperparameter tuning: find best k",14,True,WHITE)
box(s,0.5,4.15,5.8,1.0,"for k in range(1, 16):\n    5-fold CV on X_train\n    record mean accuracy\nbest_k = k with highest CV accuracy",12,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
box(s,7.0,1.15,5.8,0.4,"5-Fold Cross-Validation explained",14,True,WHITE)
box(s,7.0,1.6,5.8,2.5,"5-fold CV splits training data into 5 parts:\n\n  Fold 1: [Val] [Tr] [Tr] [Tr] [Tr]\n  Fold 2: [Tr] [Val] [Tr] [Tr] [Tr]\n  Fold 3: [Tr] [Tr] [Val] [Tr] [Tr]\n  Fold 4: [Tr] [Tr] [Tr] [Val] [Tr]\n  Fold 5: [Tr] [Tr] [Tr] [Tr] [Val]\n\nAverage accuracy across 5 folds = CV score\nTest data is NEVER seen during this step.",12,color=WHITE)

# ── Slide 9: Train CART ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Step 6 \u2014 Train CART Model")
box(s,0.5,1.15,5.8,0.4,"How CART works",14,True,WHITE)
box(s,0.5,1.6,5.8,2.1,"Builds a binary decision tree by repeatedly\nsplitting on the feature that most reduces\nGINI IMPURITY:\n\n  Gini = 1 - \u03a3(p\u1d62)\u00b2\n\n  \u2022 Gini = 0   \u2192 perfectly pure node\n  \u2022 Gini = 0.9 \u2192 very mixed node\n\nAt each node: 'Is token X score > 0.3?'\n\u2192 YES / NO branch \u2192 leaf = domain",12,color=WHITE)
box(s,0.5,3.8,5.8,0.4,"GridSearchCV parameter grid",14,True,WHITE)
box(s,0.5,4.25,5.8,0.9,"max_depth: [3, 5, 7, 10, None]\nmin_samples_split: [2, 4, 6]\n\u2192 5 \u00d7 3 = 15 combinations \u00d7 5-fold CV = 75 fits",12,color=LIGHT,bg_color=RGBColor(0x20,0x20,0x38))
box(s,7.0,1.15,5.8,0.4,"Key hyperparameters",14,True,WHITE)
for i,(p,d) in enumerate([("max_depth","Limits tree depth \u2192 prevents overfitting\nNone = grow until pure leaves"),("min_samples_split","Min samples needed to split a node\nHigher value = simpler tree"),("criterion='gini'","Splitting criterion\nGini impurity measures node purity")]):
    y=1.6+i*0.9; rect(s,7.0,y,5.8,0.82,RGBColor(0x25,0x25,0x45))
    box(s,7.1,y+0.04,5.6,0.3,p,12,True,ORANGE); box(s,7.1,y+0.38,5.6,0.42,d,11,color=WHITE)

# ── Slide 10: Evaluation Metrics ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Steps 7 & 8 \u2014 Evaluation Metrics Explained")
for i,(m,d) in enumerate([("Accuracy","% of all predictions correct\n(TP+TN) / Total"),("Precision","Of predicted domain X,\nhow many were actually X?"),("Recall","Of actual domain X,\nhow many were found?"),("F1-Score","Harmonic mean of\nPrecision & Recall")]):
    x=0.5+i*3.1; rect(s,x,1.2,2.9,1.3,RGBColor(0x25,0x25,0x45))
    box(s,x,1.22,2.9,0.45,m,15,True,ACCENT,align=PP_ALIGN.CENTER)
    box(s,x,1.7,2.9,0.75,d,11,color=WHITE,align=PP_ALIGN.CENTER)
box(s,0.5,2.75,5.8,0.4,"Confusion Matrix",14,True,WHITE)
box(s,0.5,3.2,5.8,1.6,"A 10\u00d710 matrix showing:\n  \u2022 Rows = actual domain\n  \u2022 Cols = predicted domain\n  \u2022 Diagonal = correct predictions\n  \u2022 Off-diagonal = misclassifications\nHelps identify which domains get confused.",12,color=WHITE)
box(s,7.0,2.75,5.8,0.4,"Weighted averaging",14,True,WHITE)
box(s,7.0,3.2,5.8,1.6,"Metrics computed per-class then averaged\nweighted by class size.\n\nImportant because domains are slightly\nimbalanced (SC=16 vs VS=34).\nWeighted avg prevents larger domains\nfrom masking poor performance on smaller ones.",12,color=WHITE)

# ── Slide 11: KNN Results & Interpretation *** NEW *** ───────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"KNN Results & Interpretation")
for i,(label,val,col) in enumerate([("Accuracy","83.8%",ACCENT),("Precision","84.2%",GREEN),("Recall","83.8%",ORANGE),("F1-Score","83.5%",PURPLE)]):
    stat_tile(s,0.5+i*3.05,1.15,2.9,1.15,val,label,col)
box(s,0.5,2.55,12,0.4,"What the KNN confusion matrix tells us:",14,True,WHITE)
obs=[("\u2022 Strong performers","DM, VS, LB and EX were classified near-perfectly. These domains have highly distinctive\nvariable labels (demographics, vital signs, lab tests) making them easy to separate."),
     ("\u2022 Common confusions","PE (Physical Exam) was confused with VS (Vital Signs) \u2014 both contain body measurement\nterminology, making text-based separation harder for a distance-based model."),
     ("\u2022 SC domain challenge","SC (Subject Characteristics) has only 16 training samples \u2014 the smallest class.\nKNN struggles here because fewer neighbors exist to vote correctly."),
     ("\u2022 Key limitation","KNN is a black-box: we cannot explain WHY a prediction was made, only that\nsimilar training examples voted for that domain \u2014 a problem in regulated settings.")]
y=3.05
for title,desc in obs:
    rect(s,0.5,y,12.0,0.88,RGBColor(0x25,0x25,0x45))
    box(s,0.6,y+0.04,2.5,0.3,title,11,True,ACCENT)
    box(s,0.6,y+0.38,11.5,0.44,desc,10,color=WHITE); y+=0.96

# ── Slide 12: CART Results & Interpretation *** NEW *** ──────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"CART Results & Interpretation")
for i,(label,val,col) in enumerate([("Accuracy","98.5%",GREEN),("Precision","98.6%",ACCENT),("Recall","98.5%",ORANGE),("F1-Score","98.5%",PURPLE)]):
    stat_tile(s,0.5+i*3.05,1.15,2.9,1.15,val,label,col)
box(s,0.5,2.55,12,0.4,"What the CART confusion matrix tells us:",14,True,WHITE)
obs=[("\u2022 Near-perfect accuracy","CART misclassified only 1\u20132 samples across all 10 domains on the test set.\nThe decision tree learned strong, generalizable rules from TF-IDF features."),
     ("\u2022 Top discriminating features","Tokens like 'adverse', 'laboratory', 'vital signs', 'medication' rank highest in Gini importance.\nThese mirror real-world clinical naming conventions \u2014 the model learned domain logic."),
     ("\u2022 Tree depth insight","Best max_depth from GridSearch was shallow (3\u20135 levels), meaning the classification problem\nis largely solved by a handful of high-TF-IDF tokens \u2014 the tree is not overfitting."),
     ("\u2022 Interpretability advantage","Every prediction has a clear audit trail: 'if adverse_event > 0.31 \u2192 AE'.\nCritical for clinical/regulatory environments where decisions must be explainable.")]
y=3.05
for title,desc in obs:
    rect(s,0.5,y,12.0,0.88,RGBColor(0x25,0x25,0x45))
    box(s,0.6,y+0.04,2.5,0.3,title,11,True,GREEN)
    box(s,0.6,y+0.38,11.5,0.44,desc,10,color=WHITE); y+=0.96

# ── Slide 13: Model Comparison *** NEW *** ───────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Model Comparison \u2014 KNN vs CART")
col_w=[3.5,2.5,2.5,2.5]; col_x=[0.5,4.1,6.7,9.3]
for i,h in enumerate(["Metric","KNN","CART","Winner"]):
    rect(s,col_x[i],1.2,col_w[i],0.45,ACCENT)
    box(s,col_x[i],1.2,col_w[i],0.45,h,12,True,WHITE,align=PP_ALIGN.CENTER)
rows=[("Accuracy","83.8%","98.5%","CART \u2714"),("Precision","84.2%","98.6%","CART \u2714"),
      ("Recall","83.8%","98.5%","CART \u2714"),("F1-Score","83.5%","98.5%","CART \u2714"),
      ("Interpretable?","No","Yes","CART \u2714"),("Training speed","Fast","Fast","Tie"),
      ("Prediction speed","Slow (all pairs)","Fast (tree walk)","CART \u2714")]
y=1.7
for i,row in enumerate(rows):
    bg_c=RGBColor(0x25,0x25,0x45) if i%2==0 else RGBColor(0x1E,0x1E,0x38)
    for j,cell in enumerate(row):
        rect(s,col_x[j],y,col_w[j],0.42,bg_c)
        c=GREEN if "CART" in cell else (ORANGE if "Tie" in cell else WHITE)
        box(s,col_x[j],y,col_w[j],0.42,cell,11,color=c,align=PP_ALIGN.CENTER)
    y+=0.44
box(s,0.5,5.0,12,0.4,"Verdict:",14,True,ACCENT)
box(s,0.5,5.45,12,0.75,"CART is the clear winner \u2014 15 percentage points higher accuracy AND full interpretability.\nFor SDTM domain classification, CART's decision rules align naturally with clinical naming conventions.",12,color=WHITE)

# ── Slide 14: CART Interpretability ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"CART \u2014 Decision Tree & Feature Importance")
box(s,0.5,1.15,5.8,0.4,"Decision Tree Visualization",14,True,WHITE)
box(s,0.5,1.6,5.8,1.5,"Top 3 levels of the tree are plotted.\nEach node shows:\n  \u2022 Splitting feature (TF-IDF token)\n  \u2022 Threshold value\n  \u2022 Gini impurity\n  \u2022 Sample count\n  \u2022 Predicted class (colour-coded)",12,color=WHITE)
box(s,0.5,3.3,5.8,0.4,"Feature Importances",14,True,WHITE)
box(s,0.5,3.75,5.8,1.7,"importance = total Gini reduction from all\nsplits on that feature across the tree.\n\nTop discriminating tokens:\n  'adverse'    \u2192 strongly predicts AE\n  'laboratory' \u2192 strongly predicts LB\n  'vital'      \u2192 strongly predicts VS\n  'medication' \u2192 strongly predicts CM",12,color=WHITE)
box(s,7.0,1.15,5.8,0.4,"Why interpretability matters in clinical context",14,True,WHITE)
box(s,7.0,1.6,5.8,3.6,"SDTM mapping is a regulated process.\nAuditors and data managers need to understand\nWHY a variable was classified into a domain.\n\nCART provides a clear audit trail:\n  'If adverse event > 0.3 \u2192 predict AE'\n\nKNN gives no such explanation \u2014 it only\nsays 'the 3 most similar variables in\ntraining were all AE'.\n\n\u2192 CART is more trustworthy in regulated\n  clinical environments.",12,color=WHITE)

# ── Slide 15: Summary & Conclusions ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s); header(s,"Summary & Conclusions")
for i,(title,col,points) in enumerate([
    ("KNN",ACCENT,["Instance-based (lazy learner)","Stores all training data","Cosine distance on TF-IDF","Best k chosen by 5-fold CV","No explicit decision rules","Black-box \u2014 not explainable","83.8% test accuracy"]),
    ("CART",GREEN,["Rule-based tree (eager learner)","Learns explicit split rules","Gini impurity criterion","max_depth tuned by GridSearch","Fully interpretable tree","Feature importance scores","98.5% test accuracy \u2714"])]):
    x=0.5+i*6.3; rect(s,x,1.2,5.8,0.55,col)
    box(s,x,1.2,5.8,0.55,title,20,True,WHITE,align=PP_ALIGN.CENTER)
    for j,pt in enumerate(points):
        box(s,x+0.2,1.85+j*0.44,5.5,0.42,f"\u2022 {pt}",12,color=WHITE)
box(s,0.5,5.2,12,0.4,"Key Takeaway:",14,True,ACCENT)
box(s,0.5,5.65,12,0.85,"Both models leverage TF-IDF features from SDTM variable names & labels. CART outperforms KNN by ~15 percentage points and provides interpretable decision rules \u2014 critical for clinical/regulatory use cases.",12,color=WHITE)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"C:\dev\stevens\data_analytics\project\ML_domain_classifier-\SDTM_Presentation_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
